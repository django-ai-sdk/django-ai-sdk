from __future__ import annotations

from io import BytesIO
from pathlib import Path
from unittest.mock import MagicMock

import pytest

from django_ai_sdk.files.processors import (
    DocxFileProcessor,
    PptxFileProcessor,
    XlsxFileProcessor,
)


def _make_docx(text: str = "Hello from docx") -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(text: str = "Hello from pptx") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(headers: list[str], rows: list[list[str]]) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(headers)
    for row in rows:
        ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


OFFICE_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


class TestDocxFileProcessor:
    @pytest.fixture
    def processor(self):
        return DocxFileProcessor()

    @pytest.fixture
    def real_docx_bytes(self):
        return _make_docx("Hello from docx")

    @pytest.fixture
    def mock_puremagic_docx(self, monkeypatch):
        monkeypatch.setattr(
            "django_ai_sdk.files.processors.puremagic.from_string",
            lambda stream, mime=True: OFFICE_MIME,
        )

    async def test_is_valid_docx(self, processor, tmp_path, mock_puremagic_docx):
        file = tmp_path / "test.docx"
        file.write_bytes(_make_docx())
        assert await processor.is_valid(str(file)) is True

    async def test_is_valid_rejects_wrong_extension(
        self, processor, tmp_path, mock_puremagic_docx
    ):
        file = tmp_path / "test.ppt"
        file.write_bytes(_make_docx())
        assert await processor.is_valid(str(file)) is False

    async def test_run_from_path(self, processor, tmp_path):
        file = tmp_path / "test.docx"
        file.write_bytes(_make_docx("Hello from docx"))
        result = await processor.run(str(file))
        assert result is not None
        assert "Hello from docx" in result

    async def test_run_from_path_object(self, processor, tmp_path):
        file = tmp_path / "test.docx"
        file.write_bytes(_make_docx("Path object"))
        result = await processor.run(file)
        assert result is not None
        assert "Path object" in result

    async def test_run_from_fieldfile_like(self, processor, tmp_path):
        src = tmp_path / "source.docx"
        src.write_bytes(_make_docx("FieldFile content"))
        file_like = MagicMock()
        file_like.name = "source.docx"
        file_like.path = str(src)
        result = await processor.run(file_like)
        assert result is not None
        assert "FieldFile content" in result

    async def test_is_valid_with_path_attr(
        self, processor, tmp_path, mock_puremagic_docx
    ):
        file = tmp_path / "test.docx"
        file.write_bytes(_make_docx())
        file_like = MagicMock()
        file_like.name = "test.docx"
        file_like.path = str(file)
        assert await processor.is_valid(file_like) is True


class TestPptxFileProcessor:
    @pytest.fixture
    def processor(self):
        return PptxFileProcessor()

    @pytest.fixture
    def mock_puremagic_pptx(self, monkeypatch):
        monkeypatch.setattr(
            "django_ai_sdk.files.processors.puremagic.from_string",
            lambda stream, mime=True: OFFICE_MIME,
        )

    async def test_is_valid_pptx(self, processor, tmp_path, mock_puremagic_pptx):
        file = tmp_path / "test.pptx"
        file.write_bytes(_make_pptx("Hello from pptx"))
        assert await processor.is_valid(str(file)) is True

    async def test_is_valid_rejects_wrong_extension(
        self, processor, tmp_path, mock_puremagic_pptx
    ):
        file = tmp_path / "test.doc"
        file.write_bytes(_make_pptx("Hello from pptx"))
        assert await processor.is_valid(str(file)) is False

    async def test_run(self, processor, tmp_path):
        file = tmp_path / "test.pptx"
        file.write_bytes(_make_pptx("Slide text"))
        result = await processor.run(str(file))
        assert result is not None
        assert "Slide text" in result


class TestXlsxFileProcessor:
    @pytest.fixture
    def processor(self):
        return XlsxFileProcessor()

    @pytest.fixture
    def mock_puremagic_xlsx(self, monkeypatch):
        monkeypatch.setattr(
            "django_ai_sdk.files.processors.puremagic.from_string",
            lambda stream, mime=True: OFFICE_MIME,
        )

    async def test_is_valid_xlsx(self, processor, tmp_path, mock_puremagic_xlsx):
        file = tmp_path / "test.xlsx"
        file.write_bytes(_make_xlsx(["Name", "Age"], [["Alice", "30"], ["Bob", "25"]]))
        assert await processor.is_valid(str(file)) is True

    async def test_is_valid_rejects_wrong_extension(
        self, processor, tmp_path, mock_puremagic_xlsx
    ):
        file = tmp_path / "test.csv"
        file.write_bytes(_make_xlsx(["Name", "Age"], [["Alice", "30"]]))
        assert await processor.is_valid(str(file)) is False

    async def test_run(self, processor, tmp_path):
        file = tmp_path / "test.xlsx"
        file.write_bytes(
            _make_xlsx(["Name", "Age"], [["Alice", "30"], ["Bob", "25"]])
        )
        result = await processor.run(str(file))
        assert result is not None
        assert "Alice" in result
        assert "30" in result
        assert "Bob" in result
        assert "25" in result
        assert "Name" in result
        assert "Age" in result
