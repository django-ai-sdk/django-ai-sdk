from __future__ import annotations

import asyncio
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import IO, Any, ClassVar, Protocol

import aiofiles
import puremagic
from django.conf import settings
from django.core.files.base import File

type FileSource = str | Path | File | IO[bytes]


@lru_cache
def get_allowed_files() -> dict[str, str]:
    return getattr(settings, "AI_SDK_ALLOWED_FILES", {})


def get_file_name(file: FileSource) -> str | None:
    """Get filename"""
    if isinstance(file, (str, Path)):
        return str(file)
    elif hasattr(file, "name"):
        return str(file.name)

    return None


async def read_aio_bytes(
    file: FileSource,
) -> bytes | None:
    if hasattr(file, "path"):
        async with aiofiles.open(file.path, mode="rb") as f:
            return await f.read()

    elif isinstance(file, (str, Path)):
        async with aiofiles.open(file, mode="rb") as f:
            return await f.read()

    return None


async def read_aio_text(
    file: FileSource,
    encoding: str | None = None,
) -> str | None:
    if hasattr(file, "path"):
        async with aiofiles.open(file.path, encoding=encoding) as f:
            return await f.read()

    elif isinstance(file, (str, Path)):
        async with aiofiles.open(file, encoding=encoding) as f:
            return await f.read()

    return None


async def get_mime_type(file: FileSource) -> str | None:
    stream = await read_aio_bytes(file)

    if stream is not None:
        try:
            return puremagic.from_string(stream, mime=True)
        except puremagic.PureError:
            pass

    name = get_file_name(file)
    if name:
        extension = puremagic.ext_from_filename(name)
        allowed_files = get_allowed_files()
        if extension in allowed_files:
            return allowed_files[extension]

    return None


class FileProcessor(Protocol):
    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = ()
    step: ClassVar[str | None] = None

    async def is_valid(self, file: FileSource) -> bool:
        pass

    async def run(self, file: FileSource) -> str | None:
        pass


class BaseFileProcessor:
    """Shared MIME validation for concrete file processors."""

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = ()
    step: ClassVar[str | None] = None

    async def is_valid(self, file: FileSource) -> bool:
        mime_type = await get_mime_type(file)
        return mime_type in self.ALLOWED_MIME_TYPES

    async def run(self, file: FileSource) -> str | None:
        raise NotImplementedError


class TextFileProcessor(BaseFileProcessor):
    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "text/plain",
        "text/markdown",
        "text/x-markdown",
    )

    async def run(self, file: FileSource) -> str | None:
        return await read_aio_text(file, encoding="utf-8")


class CSVFileProcessor(TextFileProcessor):
    """Returns raw CSV string. Use CSVTransform to convert to list[dict]."""

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "text/csv",
        "text/plain",
    )


class JSONFileProcessor(TextFileProcessor):
    """Returns raw JSON string. Use JSONTransform to convert to dict/list."""

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "application/json",
        "text/json",
    )


class BaseBinaryFileProcessor(BaseFileProcessor):
    """For binary formats where magic bytes alone can't distinguish types (e.g. OOXML).

    Overrides :meth:`is_valid` to also require a matching file extension,
    since puremagic returns the same MIME for all OOXML variants.
    """

    EXTENSIONS: ClassVar[tuple[str, ...]] = ()

    async def is_valid(self, file: FileSource) -> bool:
        if not await super().is_valid(file):
            return False
        name = get_file_name(file)
        if not name:
            return False
        return any(name.lower().endswith(ext) for ext in self.EXTENSIONS)


async def _get_bytes_stream(file: FileSource) -> BytesIO | None:
    data = await read_aio_bytes(file)
    if data is None:
        return None
    return BytesIO(data)


class DocxFileProcessor(BaseBinaryFileProcessor):
    """Extracts plain text from Word (.docx) files using python-docx."""

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    EXTENSIONS: ClassVar[tuple[str, ...]] = (".docx",)

    async def run(self, file: FileSource) -> str | None:
        stream = await _get_bytes_stream(file)
        if stream is None:
            return None
        return await asyncio.to_thread(self._extract_text, stream)

    @staticmethod
    def _extract_text(stream: Any) -> str:
        # Local import: docx ships in the `files` extra
        from docx import Document

        doc = Document(stream)
        return "\n".join(p.text for p in doc.paragraphs)


class PptxFileProcessor(BaseBinaryFileProcessor):
    """Extracts plain text from PowerPoint (.pptx) files using python-pptx."""

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    EXTENSIONS: ClassVar[tuple[str, ...]] = (".pptx",)

    async def run(self, file: FileSource) -> str | None:
        stream = await _get_bytes_stream(file)
        if stream is None:
            return None
        return await asyncio.to_thread(self._extract_text, stream)

    @staticmethod
    def _extract_text(stream: Any) -> str:
        # Local import: pptx ships in the `files` extra
        from pptx import Presentation

        prs = Presentation(stream)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text)
        return "\n".join(parts)


class XlsxFileProcessor(BaseBinaryFileProcessor):
    """Extracts plain text from Excel (.xlsx) files using openpyxl."""

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    EXTENSIONS: ClassVar[tuple[str, ...]] = (".xlsx", ".xls")

    async def run(self, file: FileSource) -> str | None:
        stream = await _get_bytes_stream(file)
        if stream is None:
            return None
        return await asyncio.to_thread(self._extract_text, stream)

    @staticmethod
    def _extract_text(stream: Any) -> str:
        # Local import: openpyxl ships in the `files` extra
        import openpyxl

        wb = openpyxl.load_workbook(stream, read_only=True, data_only=True)
        rows = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                rows.append(" | ".join(str(c) if c is not None else "" for c in row))
        return "\n".join(rows)


class AnyDocFileProcessor(BaseFileProcessor):
    """Converts many document formats to Markdown using anydoc."""

    step: ClassVar[str | None] = "anydoc"

    ALLOWED_MIME_TYPES: ClassVar[tuple[str, ...]] = (
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.oasis.opendocument.text",
        "application/vnd.oasis.opendocument.spreadsheet",
        "application/vnd.oasis.opendocument.presentation",
        "application/rtf",
        "application/epub+zip",
        "text/csv",
        "application/pdf",
    )

    EXTENSIONS: ClassVar[tuple[str, ...]] = (
        ".doc",
        ".docx",
        ".docm",
        ".ppt",
        ".pps",
        ".pot",
        ".pptx",
        ".pptm",
        ".ppsx",
        ".ppsm",
        ".xls",
        ".xlsx",
        ".xlsm",
        ".xlsb",
        ".odt",
        ".ods",
        ".odp",
        ".rtf",
        ".epub",
        ".csv",
        ".pdf",
    )

    async def is_valid(self, file: FileSource) -> bool:
        """Validate by extension."""
        name = get_file_name(file)
        if not name:
            return False

        # Check if the file extension is in the allowed list
        return any(name.lower().endswith(ext) for ext in self.EXTENSIONS)

    async def run(self, file: FileSource) -> str | None:
        data = await read_aio_bytes(file)
        if data is None:
            return None
        name = get_file_name(file)
        return await asyncio.to_thread(self._extract_text, data, name)

    @staticmethod
    def _extract_text(data: bytes, name: str | None) -> str | None:
        # Local import: anydoc ships in the `files` extra
        import anydoc

        fmt: str | None = anydoc.format_from_path(name) if name else None

        try:
            return anydoc.to_markdown_bytes(data, fmt) if fmt else anydoc.to_markdown_bytes(data)
        except anydoc.ConvertError:
            return None
